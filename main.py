import os
import shutil
from pptx import Presentation

# Collect fonts from PPTX (including tables)
prs = Presentation("AIFlow-AIAgens-MCP.pptx")
fonts = set()

for slide in prs.slides:
    for shape in slide.shapes:
        # Extract text frame fonts
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
        # Extract table cell fonts
        if hasattr(shape, "table"):
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.add(run.font.name)

# Prepare font copy/match
windows_fonts_folder = r"C:\Windows\Fonts"
output_folder = os.getcwd()
fonts_folder = os.path.join(output_folder, "fonts")
os.makedirs(fonts_folder, exist_ok=True)

fonts_files = os.listdir(windows_fonts_folder)
fonts_copied = []
fonts_not_copied = []

for font_name in fonts:
    found = False
    font_search_lower = font_name.lower().replace(" ", "")
    for file in fonts_files:
        file_lower = file.lower().replace(".ttf", "").replace(".otf", "")
        # Copy if font name is contained anywhere in the filename (broad matching)
        if font_search_lower in file_lower:
            src = os.path.join(windows_fonts_folder, file)
            dst = os.path.join(fonts_folder, file)
            try:
                shutil.copy(src, dst)
                fonts_copied.append(f"{font_name}: copied as {file}")
                found = True
                # If you want ALL matches, DO NOT break here
            except Exception as e:
                fonts_not_copied.append(f"{font_name}: failed to copy ({e})")
                found = True
    if not found:
        fonts_not_copied.append(f"{font_name}: {font_name} font is not available")

# Write fonts status to txt file, separating copied & not copied fonts
with open(os.path.join(output_folder, "fonts/_used_fonts_.txt"), "w", encoding="utf-8") as f:
    f.write("Fonts used in presentation:\n")
    f.write("=" * 50 + "\n\n")
    for font in sorted(fonts):
        f.write(f"- {font}\n")
    f.write("\n" + "=" * 50 + "\n")
    f.write("Font copy status:\n")
    f.write("=" * 50 + "\n\n")
    f.write("Copied fonts:\n")
    f.write("-" * 30 + "\n")
    for status in fonts_copied:
        f.write(status + "\n")
    f.write("\nNon-copied or unavailable fonts:\n")
    f.write("-" * 30 + "\n")
    for status in fonts_not_copied:
        f.write(status + "\n")

print(f"Total fonts found in PPTX: {len(fonts)}")
print(f"Fonts: {sorted(fonts)}")
print(f"Copied fonts: {len(fonts_copied)}")
print(f"Non-copied/unavailable fonts: {len(fonts_not_copied)}")
print("Fonts info saved to fonts_used.txt")
